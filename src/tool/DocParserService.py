import logging
import re
import json
import os
import time
import httpx
import aiohttp
import requests
import shutil
import urllib
import asyncio
import sys
sys.path.append('.')
sys.path.append('./api')

from tqdm import tqdm
from collections import Counter
from typing import Union, Dict, Optional, List
from tool.BaseToolService import BasicTool, register_tool
from tool.StorageService import Storage
from api.utils.key_operator import ApiKeyCycler
from api.utils.string_operator import clean_paragraph, PARAGRAPH_SPLIT_SYMBOL
from api.utils.file_operator import read_text_from_file, df_to_md, get_file_type
from api.utils.url_operator import is_http_url, sanitize_chrome_file_path, save_url_to_local_work_dir
from poi_rpc.infra.rpc.base.ttypes import Context
from poi_rpc.reddata.service import RedDataService
from poi_rpc.reddata.dto.ttypes import MultiGetNoteDetailByOidRequest
from poi_rpc.thrift_rpc.red_rpc_util import create_thrift_client
from dotenv import load_dotenv

load_dotenv('./api/utils/keys.env')
DIRECTLLM_API_KEY_USER = os.environ.get("DIRECTLLM_API_KEY_USER", "{}")
DIRECTLLM_API_KEY_USER = json.loads(DIRECTLLM_API_KEY_USER)
JINA_API_KEY = os.environ.get("JINA_API_KEY", "")
JINA_API_KEYS = os.environ.get("JINA_API_KEYS", "{}")
JINA_API_KEYS = json.loads(JINA_API_KEYS)

PARSER_SUPPORTED_FILE_TYPES = ['pdf', 'docx', 'pptx', 'txt', 'html', 'csv', 'tsv', 'xlsx', 'xls']


@register_tool('doc_parser')
class DocParser(BasicTool):
    # Reference: https://github.com/QwenLM/Qwen-Agent/blob/main/qwen_agent/tools/simple_doc_parser.py
    name = "doc_parser"
    description_en = f"Extract information, supported formats include: {' / '.join(PARSER_SUPPORTED_FILE_TYPES)}"
    description_zh = f"提取出一个文档的内容，支持类型包括：{' / '.join(PARSER_SUPPORTED_FILE_TYPES)}"
    parameters = {
        'type': 'object',
        'properties': {
            "id": {
                "type": "string",
            },
            'url': {
                'type': 'string',
            }
        },
        'required': ['url'],
    }
    def __init__(
        self,
        download_path: str = "/mnt/tidalfs-bdsz01/dataset/llm_dataset/jinjr_data/download_cache",
        cfg: Optional[Dict] = None,
        use_zh: bool = False,
        required_count_tokens: bool = False,
        max_retries: int = 1,
        retry_delay: int = 1,
        need_recompute: bool = True,
        timeout: int = 400,
        use_jina_as_backup: bool = False,
    ):
        super().__init__(
            cfg=cfg,
            use_zh=use_zh,
            max_retries=max_retries,
            retry_delay=retry_delay,
        )
        self.download_path = download_path
        self.data_storage = Storage({"root_path": self.download_path})
        self.is_structured_doc = self.cfg.get("is_structured_doc", False)
        self.need_recompute = need_recompute
        self.timeout = timeout
        
        self.required_count_tokens = required_count_tokens
        self.use_jina_as_backup = use_jina_as_backup
    
    
    @staticmethod
    def hash_sha256(text: str):
        import hashlib
        hash_object = hashlib.sha256(text.encode())
        key = hash_object.hexdigest()
        return key
    
    
    @staticmethod
    def get_plain_doc(doc: List):
        paras = []
        for page in doc:
            for para in page['content']:
                for k, v in para.items():
                    if k in ['text', 'table', 'image']:
                        paras.append(v)
        return PARAGRAPH_SPLIT_SYMBOL.join(paras)
    
    
    @staticmethod
    def parse_word(
        file_path: str,
    ):
        from docx import Document
        doc = Document(file_path)
        
        content = []
        for para in doc.paragraphs:
            content.append({'text': para.text})
        for table in doc.tables:
            tbl = []
            for row in table.rows:
                tbl.append('|' + '|'.join([cell.text for cell in row.cells]) + '|')
            tbl = '\n'.join(tbl)
            content.append({'table': tbl})
        
        # Due to the pages in Word are not fixed, the entire document is returned as one page
        return [{'page_num': 1, 'content': content}]
    
    
    @staticmethod
    def parse_ppt(
        file_path: str,
    ):
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
        
        try:
            ppt = Presentation(file_path)
        except PackageNotFoundError as ex:
            logging.info(ex)
            return []
        doc = []
        for slide_number, slide in enumerate(ppt.slides):
            page = {'page_num': slide_number + 1, 'content': []}
            
            for shape in slide.shapes:
                if not shape.has_text_frame and not shape.has_table:
                    pass
                
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = ''.join(run.text for run in paragraph.runs)
                        paragraph_text = clean_paragraph(paragraph_text)
                        if paragraph_text.strip():
                            page['content'].append({'text': paragraph_text})
                
                if shape.has_table:
                    tbl = []
                    for row_number, row in enumerate(shape.table.rows):
                        tbl.append('|' + '|'.join([cell.text for cell in row.cells]) + '|')
                    tbl = '\n'.join(tbl)
                    page['content'].append({'table': tbl})
            doc.append(page)
        return doc
    
    
    @staticmethod
    def parse_txt(
        file_path: str
    ):
        text = read_text_from_file(file_path)
        paras = text.split(PARAGRAPH_SPLIT_SYMBOL)
        content = []
        for p in paras:
            content.append({'text': p})
        
        # Due to the pages in txt are not fixed, the entire document is returned as one page
        return [{'page_num': 1, 'content': content}]
    
    
    @staticmethod
    def parse_excel(
        file_path: str,
    ):
        import pandas as pd
        excel_file = pd.ExcelFile(file_path)
        md_tables = []
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            md_table = df_to_md(df)
            md_tables.append(f'### Sheet: {sheet_name}\n{md_table}')
        
        return [{'page_num': i + 1, 'content': [{'table': md_tables[i]}]} for i in range(len(md_tables))]
    
    
    def parse_csv(
        self,
        file_path: str
    ):
        import pandas as pd
        md_tables = []
        try:
            df = pd.read_csv(file_path, encoding_errors='replace', on_bad_lines='skip')
        except Exception as ex:
            # Directly converted from Excel
            logging.info(ex)
            return self.parse_excel(file_path)
        md_table = df_to_md(df)
        md_tables.append(md_table)  # There is only one table available
        
        return [{'page_num': i + 1, 'content': [{'table': md_tables[i]}]} for i in range(len(md_tables))]
    
    
    def parse_tsv(
        self,
        file_path: str
    ):
        import pandas as pd
        md_tables = []
        try:
            df = pd.read_csv(file_path, sep='\t', encoding_errors='replace', on_bad_lines='skip')
        except Exception as ex:
            # Directly converted from Excel
            logging.info(ex)
            return self.parse_excel(file_path)
        md_table = df_to_md(df)
        md_tables.append(md_table)  # There is only one table available
        
        return [{'page_num': i + 1, 'content': [{'table': md_tables[i]}]} for i in range(len(md_tables))]
    
    
    @staticmethod
    def parse_html_bs(
        file_path: str
    ):
        def pre_process_html(s):
            # replace multiple newlines
            s = re.sub('\n+', '\n', s)
            # replace special string
            s = s.replace("Add to Qwen's Reading List", '')
            return s
        
        try:
            from bs4 import BeautifulSoup
        except Exception:
            raise ValueError('Please install bs4 by `pip install beautifulsoup4`')
        bs_kwargs = {'features': 'lxml'}
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, **bs_kwargs)
        
        text = soup.get_text()
        
        if soup.title:
            title = str(soup.title.string)
        else:
            title = ''
        
        text = pre_process_html(text)
        paras = text.split(PARAGRAPH_SPLIT_SYMBOL)
        content = []
        for p in paras:
            p = clean_paragraph(p)
            if p.strip():
                content.append({'text': p})
        
        # The entire document is returned as one page
        return [{'page_num': 1, 'content': content, 'title': title}]
    
    
    @staticmethod
    def parse_pdf(
        pdf_path: str
    ):
        # need to add header and footer
        from pdfminer.high_level import extract_pages
        from pdfminer.layout import LTImage, LTRect, LTTextContainer
        
        def postprocess_page_content(page_content: list) -> list:
            # rm repetitive identification for table and text
            # Some documents may repeatedly recognize LTRect and LTTextContainer
            table_obj = [p['obj'] for p in page_content if 'table' in p]
            tmp = []
            for p in page_content:
                repetitive = False
                if 'text' in p:
                    for t in table_obj:
                        if t.bbox[0] <= p['obj'].bbox[0] and p['obj'].bbox[1] <= t.bbox[1] and t.bbox[2] <= p['obj'].bbox[
                            2] and p['obj'].bbox[3] <= t.bbox[3]:
                            repetitive = True
                            break
                
                if not repetitive:
                    tmp.append(p)
            page_content = tmp
            
            # merge paragraphs that have been separated by mistake
            new_page_content = []
            for p in page_content:
                if new_page_content and 'text' in new_page_content[-1] and 'text' in p and abs(
                    p.get('font-size', 12) -
                    new_page_content[-1].get('font-size', 12)) < 2 and p['obj'].height < p.get('font-size', 12) + 1:
                    # Merge those lines belonging to a paragraph
                    _p = p['text']
                    new_page_content[-1]['text'] += f' {_p}'
                    # new_page_content[-1]['font-name'] = p.get('font-name', '')
                    new_page_content[-1]['font-size'] = p.get('font-size', 12)
                else:
                    p.pop('obj')
                    new_page_content.append(p)
            for i in range(len(new_page_content)):
                if 'text' in new_page_content[i]:
                    new_page_content[i]['text'] = clean_paragraph(new_page_content[i]['text'])
            return new_page_content
        
        
        def get_font(element):
            from pdfminer.layout import LTChar, LTTextContainer
            
            fonts_list = []
            for text_line in element:
                if isinstance(text_line, LTTextContainer):
                    for character in text_line:
                        if isinstance(character, LTChar):
                            fonts_list.append((character.fontname, character.size))
            
            fonts_list = list(set(fonts_list))
            if fonts_list:
                counter = Counter(fonts_list)
                most_common_fonts = counter.most_common(1)[0][0]
                return most_common_fonts
            else:
                return []
        
        
        def extract_tables(pdf, page_num):
            table_page = pdf.pages[page_num]
            tables = table_page.extract_tables()
            return tables
        
        
        def table_converter(table):
            table_string = ''
            for row_num in range(len(table)):
                row = table[row_num]
                cleaned_row = [
                    item.replace('\n', ' ') if item is not None and '\n' in item else 'None' if item is None else item
                    for item in row
                ]
                table_string += ('|' + '|'.join(cleaned_row) + '|' + '\n')
            table_string = table_string[:-1]
            return table_string
        
        doc = []
        import pdfplumber
        pdf = pdfplumber.open(pdf_path)
        for i, page_layout in enumerate(extract_pages(pdf_path)):
            page = {'page_num': page_layout.pageid, 'content': []}
            
            elements = []
            for element in page_layout:
                elements.append(element)
            
            # Init params for table
            table_num = 0
            tables = []
            
            for element in elements:
                if isinstance(element, LTRect):
                    if not tables:
                        tables = extract_tables(pdf, i)
                    if table_num < len(tables):
                        table_string = table_converter(tables[table_num])
                        table_num += 1
                        if table_string:
                            page['content'].append({'table': table_string, 'obj': element})
                elif isinstance(element, LTTextContainer):
                    # Delete line breaks in the same paragraph
                    text = element.get_text()
                    # need to further analysis using font
                    font = get_font(element)
                    if text.strip():
                        new_content_item = {'text': text, 'obj': element}
                        if font:
                            new_content_item['font-size'] = round(font[1])
                            # new_content_item['font-name'] = font[0]
                        page['content'].append(new_content_item)
                elif isinstance(element, LTImage):
                    # need to add ocr
                    raise ValueError('Currently, extracting images is not supported!')
                else:
                    pass
            
            # merge elements
            page['content'] = postprocess_page_content(page['content'])
            doc.append(page)
        
        return doc





@register_tool("web_parser")
class WebParser(DocParser):
    name = "web_parser"
    description_en = "Get content of one webpage"
    description_zh = "抓取一个网页的内容"
    parameters = {
        'type': 'object',
        'properties': {
            "web_type": {
                "type": "string",
            },
            "note_id": {
                "type": "string",
            },
            'web_url': {
                'type': 'string',
            }
        },
        'required': ['url'],
    }
    def __init__(
        self,
        download_path: str = "/mnt/tidalfs-bdsz01/dataset/llm_dataset/jinjr_data/download_cache",
        min_parsed_len: int = -1,
        max_parsed_len: int = 20480,
        cfg: Optional[Dict] = None,
        use_zh: bool = False,
        max_retries: int = 3,
        retry_delay: int = 1,
        use_jina_as_backup: bool = False,
    ):
        super().__init__(
            download_path=download_path,
            cfg=cfg,
            use_zh=use_zh,
            max_retries=max_retries,
            retry_delay=retry_delay,
            use_jina_as_backup=use_jina_as_backup,
        )
        self.min_parsed_len = min_parsed_len
        self.max_parsed_len = max_parsed_len
        self.cycler = ApiKeyCycler(api_key_list=list(JINA_API_KEYS.values()))
        
    
    async def call(
        self,
        params: Union[str, Dict],
    ):
        params = self.verify_json_format_args(params)
        url = params.get("url", "")
        
        # parse web
        note_pattern = r'/explore/([a-zA-Z0-9]+)'
        is_match = re.search(note_pattern, url)
        if is_match:
            web_type = "note"
            note_id = is_match.group(1)
        else:
            web_type = "web"
            note_id = ""
        
        parse_st = time.time()
        if not self.use_jina_as_backup:
            if web_type == "web":
                parsed_web = await self.parse_web(url=url)
            elif web_type == "note":
                parsed_web = await self.parse_note(note_id=note_id)
            else:
                parsed_web = ""
        else:
            if web_type == "note":
                parsed_web = await self.parse_note(note_id=note_id)
            else:
                parsed_web = ""
        
        parse_et = time.time()
        logging.info(f"parse web data costs: {parse_et-parse_st}")
        if web_type == "note" or (self.min_parsed_len >= 0 and len(parsed_web) >= self.min_parsed_len and len(parsed_web) > 0):
            return parsed_web[:self.max_parsed_len]
        else: # DEBUG
            web_content = await self.parse_web(url=url)
            return web_content[:self.max_parsed_len]
        
        cached_name_ori = f"{self.hash_sha256(url)}_ori"
        parsed_file = self.data_storage.get(cached_name_ori)

        if (len(parsed_file) > 0) and (not self.need_recompute):
            parsed_file = json.loads(parsed_file)
            logging.info(f'Read parsed {url} from cache.')
        
        else:
            logging.info(f'Start parsing {url} ...')
            st = time.time()
            
            f_type = get_file_type(url )
            if f_type in PARSER_SUPPORTED_FILE_TYPES:
                if url.startswith('https://') or url.startswith('http://') or re.match(
                    r'^[A-Za-z]:\\', url) or re.match(r'^[A-Za-z]:/', url):
                    url = url
                else:
                    url = sanitize_chrome_file_path(url)

            os.makedirs(self.download_path, exist_ok=True)
            if is_http_url(url):
                st = time.time()
                # download online url
                tmp_file_root = os.path.join(self.download_path, self.hash_sha256(url))
                os.makedirs(tmp_file_root, exist_ok=True)
                try:
                    url = save_url_to_local_work_dir(url, tmp_file_root)
                except:
                    url = ""
                et = time.time()

            try:
                if f_type == 'pdf':
                    parsed_file = self.parse_pdf(url)
                elif f_type == 'docx':
                    parsed_file = self.parse_word(url)
                elif f_type == 'pptx':
                    parsed_file = self.parse_ppt(url)
                elif f_type == 'txt':
                    parsed_file = self.parse_txt(url)
                elif f_type == 'html':
                    parsed_file = self.parse_html_bs(url)
                elif f_type == 'csv':
                    parsed_file = self.parse_csv(url)
                elif f_type == 'tsv':
                    parsed_file = self.parse_tsv(url)
                elif f_type in ['xlsx', 'xls']:
                    parsed_file = self.parse_excel(url)
                else:
                    _t = '/'.join(PARSER_SUPPORTED_FILE_TYPES)
                    raise ValueError(
                        f'Failed: The current parser does not support this file type! Supported types: {_t}')
            except Exception as ex:
                # raise NotImplementedError(f"{ex}")
                logging.info(f"Fail to parse file: {ex}")
                if self.use_jina_as_backup:
                    parsed_file = await self.parse_web(url=url)
                else:
                    parsed_file = ""
            
            if self.required_count_tokens:
                from api.utils.tokenizer_operator import count_tokens
                for page in parsed_file:
                    for para in page['content']:
                        # add More attribute types
                        para['token'] = count_tokens(para.get('text', para.get('table')))
            et = time.time()
            logging.info(f'Finished parsing {url}. Time spent: {et - st} seconds.')
            # Cache the parsing doc
            self.data_storage.put(cached_name_ori, json.dumps(parsed_file, ensure_ascii=False, indent=2))
        
        if not self.is_structured_doc:
            parsed_data = self.get_plain_doc(parsed_file)
        else:
            parsed_data = parsed_file
        
        if len(parsed_web) < self.min_parsed_len and len(parsed_web) > 0:
            return parsed_web[:self.max_parsed_len]
        else:
            return parsed_data[:self.max_parsed_len]
    
    
    async def parse_note(
        self,
        note_id: str,
    ):
        for attempt in range(self.max_retries):
            try:
                client2 = create_thrift_client(
                    service_name="reddataservice-service-default",
                    client_class=RedDataService,
                )
                request = MultiGetNoteDetailByOidRequest([note_id])
                response = client2.multiGetNoteDetailByOid(Context(), request)
                title = response.noteDetails[note_id].title
                content = response.noteDetails[note_id].content
                return f"{title}\n{content}"
            
            except (
                aiohttp.ClientError,
                Exception,
            ) as e:
                tqdm.write(f"请求失败 (尝试 {attempt + 1}/{self.max_retries}): {str(e)}")
                
                if attempt < self.max_retries - 1:  # 如果不是最后一次尝试
                    time.sleep(self.retry_delay)  # 等待一段时间再重试
                    continue
                else:
                    tqdm.write(f"达到最大重试次数 {self.max_retries}，放弃重试")
                    return ""  # 重试多次后仍失败，返回None

    
    async def parse_web(
        self,
        url: str,
    ):
        attempt = 0
        while attempt < self.max_retries:
            try:
                api_key = await self.cycler.get_key()
                jina_url = "https://r.jina.ai/"
                proxy = "http://10.140.15.68:3128"  # 一般要带 scheme；https 也通常用 http 代理地址
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                    "X-Engine": "browser",
                    "X-Timeout": f"{self.timeout}"
                }
                data = {
                    "url": url
                }
                async with httpx.AsyncClient(timeout=self.timeout, proxy=proxy) as client:
                    response = await client.post(jina_url, headers=headers, json=data)
                    response.raise_for_status()  # 抛出HTTP错误
                    return response.text
            except Exception as e:
                error_str = str(e)
                
                # 请求失败 (尝试 1/3): Server error '524 <none>' for url 'https://r.jina.ai/'
                # 请求失败 (尝试 1/3): Client error '451 Unavailable For Legal Reasons' for url 'https://r.jina.ai/'
                if "Server error" in error_str or "451" in error_str: # Server Error (5xx) 或 451：直接返回空
                    tqdm.write(f"不可恢复错误，直接放弃: {error_str}")
                    return ""
                
                if ("402" in error_str or
                    "Payment Required" in error_str or
                    "exceed quota" in error_str or
                    "payment" in error_str.lower()):
                    tqdm.write(f"支付或配额错误 (不计入重试次数): {error_str}")
                    await asyncio.sleep(self.retry_delay * 3)
                    continue

                # 其他错误计入重试次数
                attempt += 1
                tqdm.write(f"请求失败 (尝试 {attempt}/{self.max_retries}): {error_str}")
                
                if attempt < self.max_retries:
                    await asyncio.sleep(self.retry_delay)
                else:
                    tqdm.write(f"达到最大重试次数 {self.max_retries}，放弃重试")
                    return ""
        return ""




if __name__ == "__main__":
    async def main():
        service = WebParser(
            use_zh=True,
            max_retries=5,
            retry_delay=3,
            use_jina_as_backup=True
        )
        # params = {
        #     # "url": "https://www.xiaohongshu.com/explore/6911e587000000000703b47f",
        #     "url": "https://www.xiaohongshu.com/explore/692ce100000000001e0240c6?xsec_token=ABrGzhZDL8U9H5TM2zriiOa_dewGiMp1qRu6CEDaoKc_s=&xsec_source=pc_feed"
        # }
        params = {
            # "url": "https://proceedings.mlr.press/v80/yang18d/yang18d.pdf",
            # "url": "https://zhuanlan.zhihu.com/p/1918573796782240098"
            # "url": "https://news.sciencenet.cn/htmlnews/2025/9/551489.shtm"
            # "url": "https://www.youtube.com/watch?v=FwOTs4UxQS4"
            # "url": "http://vjs.zencdn.net/v/oceans.mp4"
            # "url": "https://cds.chinadaily.com.cn/dams/capital/image/202503/10/67ce5629e4b0ccc959925a1e_m.jpeg"
            # "url": "https://cds.chinadaily.com.cn/dams/capital/image/202503/10/67ce5629e4b0ccc959925a1e_m.jpeg"
            # "url": "http://sociology.cssn.cn/shxsw/swx_kycg/swx_xslw/202004/W020200408470418299166.pdf"
            # "url": "https://sthj.lf.gov.cn/M/Open/detail/id/124398.html"
            "url": "https://lf.hebccw.cn/system/2026/02/06/102151729.shtml"
        }
        results = await service.call(
            params=params
        )

        # results = await service.parse_web(
        #     url="https://www.xiaohongshu.com/explore/692ce100000000001e0240c6?xsec_token=ABrGzhZDL8U9H5TM2zriiOa_dewGiMp1qRu6CEDaoKc_s=&xsec_source=pc_feed"
        # )
        breakpoint()
        print(results)
    
    asyncio.run(main())
        
        
    